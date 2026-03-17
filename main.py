from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RgbColor

def create_presentation_with_drawings():
    prs = Presentation()

    # --- ДОПОМІЖНІ ФУНКЦІЇ ---

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_slide_with_text(title, lines):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        for line in lines:
            p = body.add_paragraph()
            p.text = line
            p.font.size = Pt(24)

    # Функція для малювання провідника (сірий брусок)
    def draw_conductor_scheme(slide):
        # Малюємо брусок провідника
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2), Inches(4), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(192, 192, 192) # Сірий колір
        shape.line.color.rgb = RgbColor(0, 0, 0)
        
        # Додаємо підписи зарядів на поверхні
        # Мінуси зліва
        txt_left = slide.shapes.add_textbox(Inches(4.0), Inches(2.5), Inches(0.5), Inches(0.5))
        txt_left.text_frame.paragraphs[0].text = "− − −"
        txt_left.text_frame.paragraphs[0].font.size = Pt(20)
        txt_left.text_frame.paragraphs[0].font.bold = True
        
        # Плюси справа
        txt_right = slide.shapes.add_textbox(Inches(8.2), Inches(2.5), Inches(0.5), Inches(0.5))
        txt_right.text_frame.paragraphs[0].text = "+ + +"
        txt_right.text_frame.paragraphs[0].font.size = Pt(20)
        txt_right.text_frame.paragraphs[0].font.bold = True

        # Підпис всередині "E=0"
        txt_center = slide.shapes.add_textbox(Inches(5.5), Inches(2.7), Inches(2), Inches(0.5))
        tf = txt_center.text_frame.paragraphs[0]
        tf.text = "E = 0"
        tf.font.size = Pt(32)
        tf.font.bold = True
        tf.font.color.rgb = RgbColor(255, 0, 0)

    # Функція для малювання діелектрика (жовтий брусок з диполями)
    def draw_dielectric_scheme(slide):
        # Малюємо брусок діелектрика
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2), Inches(4), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(255, 255, 150) # Світло-жовтий
        shape.line.color.rgb = RgbColor(0, 0, 0)

        # Малюємо диполі (зв'язані заряди) всередині
        # Просто додаємо текстові бокси з "+" і "-" парами
        coords = [(5.0, 2.5), (6.5, 2.5), (5.5, 3.2), (7.0, 3.2)]
        for x, y in coords:
            txt = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.8), Inches(0.5))
            txt.text_frame.paragraphs[0].text = "− +"
            txt.text_frame.paragraphs[0].font.size = Pt(22)
            txt.text_frame.paragraphs[0].font.bold = True
        
        # Стрілка поля
        txt_arrow = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2), Inches(0.5))
        txt_arrow.text_frame.paragraphs[0].text = "E ↓ (слабшає)"
        txt_arrow.text_frame.paragraphs[0].font.size = Pt(20)

    # --- СТВОРЕННЯ СЛАЙДІВ ---

    # Слайд 1: Титул
    add_title_slide("Провідники та діелектрики", "Електричне поле в речовині")

    # Слайд 2: Вступ
    add_slide_with_text(
        "Типи речовин",
        [
            "• Провідники (метали) — є вільні носії заряду.",
            "• Діелектрики (ізолятори) — немає вільних носіїв.",
            "• Напівпровідники — щось середнє."
        ]
    )

    # Слайд 3: Провідники (Зі схемою)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Провідники в електричному полі"
    
    # Текст зліва
    body = slide.shapes.placeholders[1].text_frame
    body.paragraphs[0].text = "Вільні електрони рухаються до поверхні."
    body.paragraphs[0].font.size = Pt(22)
    p1 = body.add_paragraph()
    p1.text = "Всередині поля НEMAЄ."
    p1.font.size = Pt(22)
    
    # Малюємо схему справа
    draw_conductor_scheme(slide)

    # Слайд 4: Діелектрики (Зі схемою)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Діелектрики в електричному полі"
    
    # Текст зліва
    body = slide.shapes.placeholders[1].text_frame
    body.paragraphs[0].text = "Заряди зв'язані (диполі)."
    body.paragraphs[0].font.size = Pt(22)
    p1 = body.add_paragraph()
    p1.text = "Поляризація речовини."
    p1.font.size = Pt(22)
    p2 = body.add_paragraph()
    p2.text = "Поле послаблюється."
    p2.font.size = Pt(22)

    # Малюємо схему справа
    draw_dielectric_scheme(slide)

    # Слайд 5: Висновок
    add_slide_with_text(
        "Висновки",
        [
            "1. Провідник екранує поле (E=0 всередині).",
            "2. Діелектрик послаблює поле (E зменшується).",
            "3. Це використовують в техніці та побуті."
        ]
    )

    prs.save('Physics_Drawings.pptx')
    print("Презентація 'Physics_Drawings.pptx' створена з намальованими схемами!")

if __name__ == "__main__":
    create_presentation_with_drawings()
